from __future__ import annotations

import io
import math
import re
import shutil
import subprocess
import tempfile
import zipfile
from collections import defaultdict
from pathlib import Path

import streamlit as st
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".webp", ".tif", ".tiff"}
PHOTO_FOLDER_NAMES = {"照片", "photo", "photos", "image", "images", "图片"}


def natural_key(value: str) -> list[object]:
    parts: list[object] = []
    chunk = ""
    is_digit = False
    for char in value:
        if char.isdigit() != is_digit and chunk:
            parts.append(int(chunk) if is_digit else chunk.casefold())
            chunk = ""
        chunk += char
        is_digit = char.isdigit()
    if chunk:
        parts.append(int(chunk) if is_digit else chunk.casefold())
    return parts


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001


def extract_archive(archive_path: Path, destination: Path) -> None:
    suffix = archive_path.suffix.lower()
    if suffix == ".zip":
        with zipfile.ZipFile(archive_path) as archive:
            extract_zip_with_chinese_names(archive, destination)
        return

    if suffix == ".rar":
        if shutil.which("bsdtar") is None:
            raise RuntimeError("当前系统没有可用的 RAR 解压工具 bsdtar。")
        subprocess.run(["bsdtar", "-xf", str(archive_path), "-C", str(destination)], check=True)
        return

    raise RuntimeError("照片包请上传 .zip 或 .rar 文件。")


def decode_zip_member_name(member: zipfile.ZipInfo) -> str:
    if member.flag_bits & 0x800:
        return member.filename

    try:
        raw_name = member.filename.encode("cp437")
    except UnicodeEncodeError:
        return member.filename

    for encoding in ("utf-8", "gb18030", "gbk"):
        try:
            return raw_name.decode(encoding)
        except UnicodeDecodeError:
            continue
    return member.filename


def safe_destination(root: Path, member_name: str) -> Path:
    normalized = Path(member_name)
    if normalized.is_absolute() or ".." in normalized.parts:
        raise RuntimeError(f"压缩包里包含不安全路径：{member_name}")

    target = root / normalized
    target_parent = target.parent.resolve()
    root_resolved = root.resolve()
    if root_resolved not in (target_parent, *target_parent.parents):
        raise RuntimeError(f"压缩包里包含不安全路径：{member_name}")
    return target


def extract_zip_with_chinese_names(archive: zipfile.ZipFile, destination: Path) -> None:
    for member in archive.infolist():
        member_name = decode_zip_member_name(member)
        target = safe_destination(destination, member_name)
        if member.is_dir():
            target.mkdir(parents=True, exist_ok=True)
            continue

        target.parent.mkdir(parents=True, exist_ok=True)
        with archive.open(member) as source, target.open("wb") as output:
            shutil.copyfileobj(source, output)


def is_hidden_or_system(path: Path) -> bool:
    return any(part.startswith(".") or part == "__MACOSX" for part in path.parts)


def find_image_groups(root: Path) -> list[tuple[str, list[Path]]]:
    grouped: dict[str, list[Path]] = defaultdict(list)
    for path in root.rglob("*"):
        if not path.is_file() or is_hidden_or_system(path.relative_to(root)):
            continue
        if path.suffix.lower() not in IMAGE_EXTENSIONS:
            continue
        parent_name = path.parent.name
        if parent_name.casefold() in PHOTO_FOLDER_NAMES:
            parent_name = path.stem
        grouped[parent_name].append(path)

    result = []
    for folder_name, images in grouped.items():
        ordered_images = sorted(images, key=lambda item: natural_key(item.name))
        result.append((folder_name, ordered_images))
    return sorted(result, key=lambda item: natural_key(item[0]))


def normalize_route_plate(folder_name: str, *, vehicle_prefix: str, auto_prefix: bool) -> str:
    if not auto_prefix:
        return folder_name

    clean_prefix = vehicle_prefix.strip()
    if not clean_prefix or clean_prefix in folder_name:
        return folder_name

    if "-" in folder_name:
        route_part, plate_part = folder_name.rsplit("-", 1)
        route_prefix = f"{route_part}-"
    else:
        route_prefix = ""
        plate_part = folder_name

    match = re.search(r"([A-Za-z0-9]{5,8})$", plate_part)
    if match is None:
        return folder_name

    plate_tail = match.group(1).upper()
    if clean_prefix[-1:].isascii() and plate_tail.startswith(clean_prefix[-1:].upper()) and len(plate_tail) > 5:
        plate_tail = plate_tail[1:]

    return f"{route_prefix}{clean_prefix}{plate_tail}"


def add_text(cell, text: str, *, bold: bool = False, size: int = 14) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cell.text = ""
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = RGBColor(0, 0, 0)
    r_pr = run._r.get_or_add_rPr()  # noqa: SLF001 - required for East Asian font hints.
    for font_tag in ("a:latin", "a:ea", "a:cs"):
        font_element = r_pr.find(qn(font_tag))
        if font_element is None:
            font_element = parse_xml(f'<{font_tag} {nsdecls("a")}/>')
            r_pr.append(font_element)
        font_element.set("typeface", "Microsoft YaHei")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_border(cell, color: str = "000000", width: str = "12700") -> None:
    tc = cell._tc  # noqa: SLF001 - python-pptx exposes table borders through XML only.
    tc_pr = tc.get_or_add_tcPr()
    for line in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        element = tc_pr.find(line, namespaces=tc_pr.nsmap)
        if element is None:
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import nsdecls

            element = parse_xml(f'<{line} {nsdecls("a")} w="{width}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></{line}>')
            tc_pr.append(element)


def add_info_table(slide, *, customer: str, brand: str, route_plate: str, publish_form: str, page_width: int) -> None:
    margin = page_width * 0.065
    table_width = page_width - (margin * 2)
    table_height = Inches(0.82)
    table = slide.shapes.add_table(2, 4, margin, Inches(0.48), table_width, table_height).table

    col_widths = [table_width * 0.145, table_width * 0.455, table_width * 0.145, table_width * 0.255]
    for index, width in enumerate(col_widths):
        table.columns[index].width = int(width)

    values = [
        ("客户", customer, "品牌", brand),
        ("线路车牌", route_plate, "发布形式", publish_form),
    ]
    for row_index, row_values in enumerate(values):
        for col_index, text in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            is_label = col_index in (0, 2)
            add_text(cell, text, bold=is_label, size=12 if is_label else 11)
            set_cell_border(cell)


def open_display_image(image_path: Path) -> Image.Image:
    with Image.open(image_path) as image:
        display_image = ImageOps.exif_transpose(image)
        display_image.load()
    if display_image.mode in ("RGBA", "LA"):
        background = Image.new("RGB", display_image.size, "white")
        alpha = display_image.getchannel("A")
        background.paste(display_image.convert("RGB"), mask=alpha)
        return background
    return display_image.convert("RGB")


def image_as_jpeg_stream(image_path: Path, *, max_side: int = 1800, quality: int = 85) -> io.BytesIO:
    image = open_display_image(image_path)
    try:
        image.thumbnail((max_side, max_side), Image.Resampling.LANCZOS)
    except AttributeError:
        image.thumbnail((max_side, max_side), Image.LANCZOS)
    stream = io.BytesIO()
    image.save(stream, format="JPEG", quality=quality, optimize=True)
    stream.seek(0)
    image.close()
    return stream


def get_fit_geometry(image_path: Path, frame_left: int, frame_top: int, frame_width: int, frame_height: int) -> tuple[int, int, int, int]:
    image = open_display_image(image_path)
    try:
        width, height = image.size
    finally:
        image.close()
    image_ratio = width / height
    frame_ratio = frame_width / frame_height
    if image_ratio >= frame_ratio:
        draw_width = frame_width
        draw_height = int(frame_width / image_ratio)
    else:
        draw_height = frame_height
        draw_width = int(frame_height * image_ratio)
    left = frame_left + int((frame_width - draw_width) / 2)
    top = frame_top + int((frame_height - draw_height) / 2)
    return left, top, draw_width, draw_height


def add_images(
    slide,
    images: list[Path],
    *,
    page_width: int,
    page_height: int,
    max_image_side: int,
    image_quality: int,
) -> None:
    left = int(page_width * 0.16)
    width = int(page_width * 0.68)
    top = Inches(1.55)
    bottom_margin = Inches(0.3)
    gap = Inches(0.32)
    available_height = page_height - top - bottom_margin
    slot_count = max(1, min(2, len(images)))
    slot_height = int((available_height - gap * (slot_count - 1)) / slot_count)

    for index, image_path in enumerate(images[:2]):
        frame_top = int(top + index * (slot_height + gap))
        img_left, img_top, img_width, img_height = get_fit_geometry(image_path, left, frame_top, width, slot_height)
        slide.shapes.add_picture(
            image_as_jpeg_stream(image_path, max_side=max_image_side, quality=image_quality),
            img_left,
            img_top,
            width=img_width,
            height=img_height,
        )


def blank_slide_layout(prs: Presentation):
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def build_ppt(
    *,
    template_path: Path | None,
    groups: list[tuple[str, list[Path]]],
    customer: str,
    brand: str,
    publish_form: str,
    vehicle_prefix: str = "辽B",
    auto_vehicle_prefix: bool = True,
    max_image_side: int = 1800,
    image_quality: int = 85,
    output_path: Path,
) -> int:
    prs = Presentation(str(template_path)) if template_path else Presentation()
    remove_all_slides(prs)
    layout = blank_slide_layout(prs)

    page_count = 0
    for route_plate, images in groups:
        display_route_plate = normalize_route_plate(
            route_plate,
            vehicle_prefix=vehicle_prefix,
            auto_prefix=auto_vehicle_prefix,
        )
        page_total = max(1, math.ceil(len(images) / 2))
        for page_index in range(page_total):
            page_images = images[page_index * 2 : page_index * 2 + 2]
            slide = prs.slides.add_slide(layout)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            add_info_table(
                slide,
                customer=customer,
                brand=brand,
                route_plate=display_route_plate,
                publish_form=publish_form,
                page_width=prs.slide_width,
            )
            add_images(
                slide,
                page_images,
                page_width=prs.slide_width,
                page_height=prs.slide_height,
                max_image_side=max_image_side,
                image_quality=image_quality,
            )
            page_count += 1

    prs.save(output_path)
    return page_count


def save_uploaded_file(uploaded_file, destination: Path) -> Path:
    destination.write_bytes(uploaded_file.getbuffer())
    return destination


def render_app(*, configure_page: bool = False) -> None:
    if configure_page:
        st.set_page_config(page_title="车体监测 PPT 自动生成", layout="wide")
    st.title("车体监测 PPT 自动生成")

    left, right = st.columns([0.9, 1.1], gap="large")
    with left:
        template_file = st.file_uploader("上传 PPT 模版", type=["pptx"])
        photo_archive = st.file_uploader("上传照片压缩包", type=["zip", "rar"])
        customer = st.text_input("客户", value="利洁时（中国）投资有限公司")
        brand = st.text_input("品牌", value="杜蕾斯")
        publish_form = st.text_input("发布形式", value="大三侧")
        auto_vehicle_prefix = st.checkbox("自动补车牌前缀", value=True)
        vehicle_prefix = st.text_input("车牌前缀", value="辽B", disabled=not auto_vehicle_prefix)
        with st.expander("图片压缩设置"):
            max_image_side = st.slider("图片最大边长", min_value=1000, max_value=2600, value=1800, step=100)
            image_quality = st.slider("图片质量", min_value=70, max_value=95, value=85, step=1)
        output_name = st.text_input("导出文件名", value="车体监测报告.pptx")

    with right:
        st.markdown(
            """
            **照片包结构示例**

            ```text
            照片/
              6-0405-辽BM2910/
                1.jpg
                2.jpg
                3.jpg
                4.jpg
              10-6537-辽B09691D/
                1.jpg
                2.jpg
            ```

            每个文件夹名会自动写入“线路车牌”。同一文件夹每 2 张照片生成 1 页，超过 2 张会继续生成下一页，表格内容保持一致。
            """
        )

    can_generate = template_file is not None and photo_archive is not None
    if st.button("生成 PPT", type="primary", disabled=not can_generate):
        if not output_name.lower().endswith(".pptx"):
            output_name = f"{output_name}.pptx"

        with st.spinner("正在解压照片并生成 PPT..."):
            with tempfile.TemporaryDirectory() as temp_dir_name:
                temp_dir = Path(temp_dir_name)
                template_path = save_uploaded_file(template_file, temp_dir / template_file.name)
                archive_path = save_uploaded_file(photo_archive, temp_dir / photo_archive.name)
                extract_dir = temp_dir / "photos"
                extract_dir.mkdir()
                extract_archive(archive_path, extract_dir)
                groups = find_image_groups(extract_dir)

                if not groups:
                    st.error("没有在压缩包里找到图片文件。")
                    return

                output_path = temp_dir / output_name
                page_count = build_ppt(
                    template_path=template_path,
                    groups=groups,
                    customer=customer,
                    brand=brand,
                    publish_form=publish_form,
                    vehicle_prefix=vehicle_prefix,
                    auto_vehicle_prefix=auto_vehicle_prefix,
                    max_image_side=max_image_side,
                    image_quality=image_quality,
                    output_path=output_path,
                )
                ppt_bytes = output_path.read_bytes()

        st.success(f"已生成 {page_count} 页，包含 {len(groups)} 个线路车牌。")
        st.download_button(
            "下载生成的 PPT",
            data=ppt_bytes,
            file_name=output_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


if __name__ == "__main__":
    render_app(configure_page=True)
